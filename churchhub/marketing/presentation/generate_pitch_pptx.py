#!/usr/bin/env python3
"""Generate ChurchHub Enterprise pitch deck (PPTX)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "ChurchHub_Enterprise_Pitch.pptx"

NAVY = RGBColor(0x0B, 0x1F, 0x3A)
BLUE = RGBColor(0x1F, 0x6F, 0xEB)
SLATE = RGBColor(0x33, 0x41, 0x55)
LIGHT = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x64, 0x74, 0x8B)


def _set_run_font(run, size_pt, bold=False, color=NAVY, name="Calibri"):
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def _fill_solid(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def _add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run_font(run, size, bold=bold, color=color)
    return box


def _add_bullets(slide, left, top, width, height, lines, size=16, color=SLATE):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = line
        _set_run_font(run, size, color=color)
    return box


def _add_notes(slide, text: str):
    notes = slide.notes_slide
    notes.notes_text_frame.text = text.strip()


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def _banner(slide, title: str, subtitle: str = ""):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.05))
    _fill_solid(bar, NAVY)
    bar.line.fill.background()
    _add_textbox(slide, Inches(0.5), Inches(0.22), Inches(12), Inches(0.5), title, size=28, bold=True, color=WHITE)
    if subtitle:
        _add_textbox(slide, Inches(0.5), Inches(0.65), Inches(12), Inches(0.35), subtitle, size=12, color=RGBColor(0x93, 0xC5, 0xFD))


def _footer(slide, page: str):
    _add_textbox(
        slide,
        Inches(0.5),
        Inches(7.15),
        Inches(10),
        Inches(0.3),
        "ChurchHub Enterprise  ·  Sample demo data where charts appear  ·  Confidential",
        size=10,
        color=MUTED,
    )
    _add_textbox(slide, Inches(11.5), Inches(7.15), Inches(1.3), Inches(0.3), page, size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def _chip(slide, left, top, text, width=1.7):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(width), Inches(0.42))
    _fill_solid(shape, LIGHT)
    shape.line.color.rgb = BLUE
    tf = shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = text
    _set_run_font(run, 11, bold=True, color=NAVY)
    tf.word_wrap = False


def _icon_label(slide, left, top, glyph: str, label: str):
    """Unicode glyph stands in for Bootstrap Icons until SVGs are dropped in."""
    _add_textbox(slide, left, top, Inches(1.8), Inches(0.4), glyph, size=28, color=BLUE, align=PP_ALIGN.CENTER)
    _add_textbox(slide, left, top + Inches(0.4), Inches(1.8), Inches(0.4), label, size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- 01 Title ---
    s = _blank_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    _fill_solid(bg, NAVY)
    bg.line.fill.background()
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    _fill_solid(accent, BLUE)
    accent.line.fill.background()
    _add_textbox(s, Inches(0.9), Inches(2.0), Inches(11), Inches(0.8), "CHURCHHUB", size=54, bold=True, color=WHITE)
    _add_textbox(s, Inches(0.9), Inches(2.85), Inches(11), Inches(0.45), "Enterprise Church Management System", size=22, color=RGBColor(0x93, 0xC5, 0xFD))
    _add_textbox(s, Inches(0.9), Inches(3.6), Inches(11), Inches(0.5), "One platform. Complete hierarchy. Books you can trust.", size=20, color=WHITE)
    _add_textbox(s, Inches(0.9), Inches(6.5), Inches(11), Inches(0.35), "Confidential  ·  Prospect briefing  ·  Icon: building / brand mark", size=12, color=MUTED)
    _add_notes(
        s,
        "Good morning. In the next ten minutes I’ll show ChurchHub — multi-level church administration "
        "from the local church to conference and beyond. Roles, church scope, and audit throughout.",
    )

    # --- 02 Agenda ---
    s = _blank_slide(prs)
    _banner(s, "Ten minutes. The full arc.", "Agenda  ·  Icons: speedometer · people · cash · chart · calendar · building-gear")
    items = [
        ("⏱", "1  Mission Control"),
        ("👥", "2  People & visitors"),
        ("💵", "3  Giving & treasury"),
        ("📊", "4  Reports"),
        ("📅", "5  Church life"),
        ("🏢", "6  Platform"),
    ]
    for i, (g, lab) in enumerate(items):
        _icon_label(s, Inches(0.6 + i * 2.1), Inches(2.4), g, lab)
    _add_textbox(s, Inches(0.5), Inches(4.5), Inches(12), Inches(1.2),
                 "Live demo preferred. These slides are the map — and the backup if the app is unavailable.",
                 size=16, color=SLATE)
    _footer(s, "02")
    _add_notes(s, "Ten minutes. The full arc: Mission Control, people, giving and treasury, reports, church life, platform.")

    # --- 03 Challenge ---
    s = _blank_slide(prs)
    _banner(s, "Church administration is enterprise work", "The challenge")
    _add_bullets(s, Inches(0.5), Inches(1.4), Inches(6.2), Inches(3.5), [
        "• Spreadsheets break at scale",
        "• Point tools don’t respect hierarchy",
        "• Finance mistakes erode trust",
        "• Pastoral follow-ups get lost",
    ], size=18)
    quote = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.5), Inches(5.5), Inches(1.8))
    _fill_solid(quote, LIGHT)
    quote.line.color.rgb = BLUE
    tf = quote.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "“Financial mistakes aren’t IT issues — they’re trust issues.”"
    _set_run_font(run, 16, bold=True, color=NAVY)

    chart_data = CategoryChartData()
    chart_data.categories = ["Spreadsheets", "Point tools", "Unified ChMS"]
    chart_data.add_series("Systems in use (sample)", (4, 3, 1))
    chart = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(7.0), Inches(3.5), Inches(5.5), Inches(3.2), chart_data).chart
    chart.has_legend = False
    _footer(s, "03")
    _add_notes(s, "Pitch-only slide. Spreadsheets and point tools fail hierarchy and trust. ChurchHub replaces fragmentation.")

    # --- 04 Product ---
    s = _blank_slide(prs)
    _banner(s, "Built for your structure", "Not a flat single-site toy  ·  Icon: diagram-3")
    levels = ["General Conference", "Union", "Conference", "Zone", "District", "Local Church"]
    for i, lvl in enumerate(levels):
        y = Inches(1.35 + i * 0.7)
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(5.8), Inches(0.55))
        _fill_solid(shape, NAVY if i % 2 == 0 else BLUE)
        shape.line.fill.background()
        tf = shape.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = lvl
        _set_run_font(run, 14, bold=True, color=WHITE)
    chips = ["Membership", "Treasury", "Remittance", "Church Life", "Reports", "Portal", "Platform"]
    for i, c in enumerate(chips):
        _chip(s, Inches(7.0 + (i % 2) * 2.8), Inches(1.5 + (i // 2) * 0.7), c, width=2.5)
    _add_textbox(s, Inches(7.0), Inches(5.5), Inches(5.5), Inches(1.0),
                 "Screenshot drop-in: 03-organization-hierarchy.png", size=12, color=MUTED)
    _footer(s, "04")
    _add_notes(s, "Hierarchy from GC to local church. One platform for membership, treasury, remittance, church life, reports, portal, platform.")

    # --- 05 Mission Control ---
    s = _blank_slide(prs)
    _banner(s, "Mission Control — pulse for every role", "Icons: speedometer2 · list-check · calculator · lightning")
    _add_bullets(s, Inches(0.5), Inches(1.35), Inches(6), Inches(3.2), [
        "• Role-aware KPIs",
        "• Action Queue (visitors, transfers, approvals)",
        "• Teller console & business date",
        "• This Week Pulse — birthdays, meetings, follow-ups",
    ], size=17)
    chart_data = CategoryChartData()
    chart_data.categories = ["Feb", "Mar", "Apr", "May", "Jun", "Jul"]
    chart_data.add_series("Income", (42, 45, 48, 51, 49, 55))
    chart_data.add_series("Expense", (28, 30, 29, 33, 31, 34))
    chart = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(6.8), Inches(1.35), Inches(5.9), Inches(4.6), chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.chart_title.text_frame.paragraphs[0].text = "Income vs Expense (sample)"
    _footer(s, "05")
    _add_notes(
        s,
        "Mission Control: KPIs scoped to active church, Action Queue, teller/business date, pastoral pulse. "
        "Chart is sample demo data.",
    )

    # --- 06 People ---
    s = _blank_slide(prs)
    _banner(s, "Membership & visitors — one living record", "Icons: person-vcard · people · person-plus · heartbeat")
    _add_textbox(s, Inches(0.5), Inches(1.3), Inches(5.5), Inches(0.4), "MEMBERS", size=14, bold=True, color=BLUE)
    _add_bullets(s, Inches(0.5), Inches(1.7), Inches(5.5), Inches(2.8), [
        "• Directory, profile, journey timeline",
        "• Families, leadership, baptisms",
        "• Transfers with audit history",
    ], size=16)
    _add_textbox(s, Inches(0.5), Inches(4.4), Inches(5.5), Inches(0.4), "VISITORS", size=14, bold=True, color=BLUE)
    _add_bullets(s, Inches(0.5), Inches(4.8), Inches(5.5), Inches(1.8), [
        "• Capture & follow-up assignment",
        "• Convert to member — history preserved",
    ], size=16)
    chart_data = CategoryChartData()
    chart_data.categories = ["Baptized", "Profession", "Visitor track", "Other"]
    chart_data.add_series("Share", (62, 18, 12, 8))
    chart = s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.0), chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    _footer(s, "06")
    _add_notes(s, "Living directory with church scope. Journey preserves history. Visitors convert without re-typing.")

    # --- 07 Stewardship ---
    s = _blank_slide(prs)
    _banner(s, "Giving you can counsel. Books you can trust.", "Icons: receipt · printer · bank · calendar3")
    _add_bullets(s, Inches(0.5), Inches(1.35), Inches(6.2), Inches(3.5), [
        "• Permissioned giving statements",
        "• Record receipt → printable confirmation",
        "• Maker-checker approvals",
        "• Business date & financial periods",
        "• Remittance / monthly cut-off",
    ], size=16)
    chart_data = CategoryChartData()
    chart_data.categories = ["Central", "Eastside", "Riverside"]
    chart_data.add_series("Tithe", (18.2, 12.1, 9.7))
    chart_data.add_series("Combined", (6.4, 4.8, 3.2))
    chart = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, Inches(6.9), Inches(1.35), Inches(5.8), Inches(4.8), chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.chart_title.text_frame.paragraphs[0].text = "Tithe vs Combined (MTD, sample $k)"
    _footer(s, "07")
    _add_notes(
        s,
        "Permissioned giving. Business-date posting. Printable confirmation. Maker-checker. Remittance cut-off at conference.",
    )

    # --- 08 Church Life ---
    s = _blank_slide(prs)
    _banner(s, "Meetings, calendar, announcements", "One pastoral rhythm  ·  Icons: calendar-week · journal-text · megaphone · phone")
    cards = [
        ("📅 Calendar", "Meetings · birthdays · communications — one shared truth"),
        ("📝 Meetings", "Agenda · minutes · decisions · attendance · approvals"),
        ("📣 + 📱 Portal", "Announcements publish to church and member phones"),
    ]
    for i, (title, body) in enumerate(cards):
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5 + i * 4.2), Inches(2.0), Inches(3.9), Inches(3.5))
        _fill_solid(card, LIGHT)
        card.line.color.rgb = BLUE
        _add_textbox(s, Inches(0.7 + i * 4.2), Inches(2.3), Inches(3.5), Inches(0.6), title, size=18, bold=True, color=NAVY)
        _add_textbox(s, Inches(0.7 + i * 4.2), Inches(3.1), Inches(3.5), Inches(2.0), body, size=15, color=SLATE)
    _add_textbox(s, Inches(0.5), Inches(5.8), Inches(12), Inches(0.5),
                 "Screenshot drop-ins: 07-announcements-calendar · meetings-detail · 08-portal-mobile", size=12, color=MUTED)
    _footer(s, "08")
    _add_notes(s, "Church Life: shared calendar, governed meetings, announcements to portal.")

    # --- 09 Reports ---
    s = _blank_slide(prs)
    _banner(s, "Reports leaders use. Integrity auditors expect.", "Icons: file-earmark-bar-graph · balance-scale · download")
    _add_bullets(s, Inches(0.5), Inches(1.35), Inches(6.5), Inches(3.8), [
        "• Report Center — role-filtered catalog",
        "• Tithe & offering · membership · attendance",
        "• Trial balance · income statement · balance sheet",
        "• Export CSV · Excel · PDF",
        "• Double-entry books of record — not a second GL",
    ], size=16)
    # Mini trial balance table as text block
    table_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.5), Inches(5.4), Inches(4.2))
    _fill_solid(table_box, LIGHT)
    table_box.line.color.rgb = NAVY
    _add_textbox(s, Inches(7.4), Inches(1.7), Inches(5), Inches(0.4), "Sample trial balance", size=14, bold=True, color=NAVY)
    _add_textbox(
        s,
        Inches(7.4),
        Inches(2.3),
        Inches(5),
        Inches(3.0),
        "Cash ................ 120,000 Dr\n"
        "Tithe Income ........ 85,000 Cr\n"
        "Offerings ........... 25,000 Cr\n"
        "Expenses ............ 15,000 Dr\n"
        "─────────────────────────\n"
        "Total  135,000 = 135,000  ✓",
        size=14,
        color=SLATE,
    )
    _footer(s, "09")
    _add_notes(s, "One report center. Stewardship answers in seconds. Trial balance equality — double-entry integrity.")

    # --- 10 Security & Platform ---
    s = _blank_slide(prs)
    _banner(s, "Least privilege. Denomination wall. Operator lane.", "Icons: key · shield-lock · layers · toggles · journal-check")
    pillars = [
        ("RBAC", "Roles, overrides,\neffective permissions"),
        ("Tenancy", "Church scope +\ndenomination isolation"),
        ("Platform", "Subscriptions, features,\naudit, health"),
    ]
    for i, (t, b) in enumerate(pillars):
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5 + i * 4.2), Inches(1.5), Inches(3.9), Inches(2.4))
        _fill_solid(card, NAVY if i == 2 else LIGHT)
        card.line.color.rgb = BLUE
        title_c = WHITE if i == 2 else NAVY
        body_c = RGBColor(0xE2, 0xE8, 0xF0) if i == 2 else SLATE
        _add_textbox(s, Inches(0.7 + i * 4.2), Inches(1.75), Inches(3.5), Inches(0.5), t, size=20, bold=True, color=title_c)
        _add_textbox(s, Inches(0.7 + i * 4.2), Inches(2.4), Inches(3.5), Inches(1.2), b, size=14, color=body_c)
    # Architecture strip
    for i, label in enumerate(["Institution Workspace", "Denomination Wall", "Platform Control"]):
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5 + i * 4.2), Inches(4.4), Inches(3.9), Inches(1.5))
        _fill_solid(shape, BLUE if i == 1 else LIGHT)
        shape.line.fill.background()
        tf = shape.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = label
        _set_run_font(run, 14, bold=True, color=WHITE if i == 1 else NAVY)
    _footer(s, "10")
    _add_notes(
        s,
        "RBAC, tenancy wall, separate platform lane. Features per church. Audited operator actions.",
    )

    # --- 11 Roles ---
    s = _blank_slide(prs)
    _banner(s, "What each leader walks away with", "Outcomes by role")
    roles = [
        ("Executive", "Visibility & hierarchy roll-up", "briefcase"),
        ("Pastor", "Action Queue & visitors", "heart"),
        ("Treasurer", "Receipts, statements, trial balance", "wallet"),
        ("Admin / Clerk", "Directory, transfers, invitations", "gear"),
    ]
    for i, (t, b, _icon) in enumerate(roles):
        col = i % 2
        row = i // 2
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5 + col * 6.4), Inches(1.5 + row * 2.5), Inches(6.0), Inches(2.2))
        _fill_solid(card, LIGHT)
        card.line.color.rgb = BLUE
        _add_textbox(s, Inches(0.8 + col * 6.4), Inches(1.75 + row * 2.5), Inches(5.4), Inches(0.5), t, size=22, bold=True, color=NAVY)
        _add_textbox(s, Inches(0.8 + col * 6.4), Inches(2.4 + row * 2.5), Inches(5.4), Inches(0.8), b, size=16, color=SLATE)
    _footer(s, "11")
    _add_notes(s, "Executives: visibility. Pastors: queue. Treasurers: books. Admins: membership and access. One system.")

    # --- 12 CTA ---
    s = _blank_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    _fill_solid(bg, NAVY)
    bg.line.fill.background()
    _add_textbox(s, Inches(0.9), Inches(1.5), Inches(11), Inches(0.6), "Recommended next step", size=18, color=RGBColor(0x93, 0xC5, 0xFD))
    _add_textbox(s, Inches(0.9), Inches(2.1), Inches(11), Inches(1.0), "Guided pilot on your conference hierarchy", size=32, bold=True, color=WHITE)
    _add_bullets(s, Inches(0.9), Inches(3.3), Inches(10), Inches(1.5), [
        "• Your chart of accounts",
        "• Two local churches",
        "• Your roles and remittance rules",
    ], size=18, color=RGBColor(0xE2, 0xE8, 0xF0))
    for i, label in enumerate(["Request a Demo", "Talk to Sales", "Start Pilot"]):
        btn = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9 + i * 3.8), Inches(5.3), Inches(3.4), Inches(0.7))
        _fill_solid(btn, BLUE)
        btn.line.fill.background()
        tf = btn.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = label
        _set_run_font(run, 16, bold=True, color=WHITE)
    _add_notes(
        s,
        "Recommended next step: guided pilot with your hierarchy, chart of accounts, and two churches. "
        "Happy to go deeper on remittance, permissions, or the member portal. What matters most to your team?",
    )

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
